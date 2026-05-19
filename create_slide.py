from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
slide_layout = prs.slide_layouts[5] # Title only
slide = prs.slides.add_slide(slide_layout)

# Title
title = slide.shapes.title
title.text = "Brand Awareness Gap in Korea"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.size = Pt(36)

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Criteo dominates local content & SEO, while RTB House lacks Korean visibility."
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(100, 100, 100)

# Add Table
rows = 6
cols = 3
left = Inches(0.5)
top = Inches(1.9)
width = Inches(9.0)
height = Inches(2.2)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(4.2)
table.columns[1].width = Inches(2.4)
table.columns[2].width = Inches(2.4)

# Headers
headers = ['Local Market Metric (Korea)', 'Criteo Korea', 'RTB House Korea']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

# Data rows
data = [
    ['Active Local Blog, Website & SEO', 'O', 'X'],
    ['Quarterly Korean Whitepapers & Reports', 'O', 'X'],
    ['Local PR, News & Executive Interviews', 'O', 'X'],
    ['Public Local Case Studies (Tier-1 Brands)', 'O', 'X'],
    ['Flagship Offline Events (e.g., Summit)', 'O', 'X']
]

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER
        
        # Color coding for O and X
        if text == 'O':
            p.font.color.rgb = RGBColor(0, 153, 0)
            p.font.bold = True
        elif text == 'X':
            p.font.color.rgb = RGBColor(204, 0, 0)
            p.font.bold = True
        else:
            # Align first column to left
            p.alignment = PP_ALIGN.LEFT

# Action Plan
txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(2.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True

p1 = tf2.add_paragraph()
p1.text = "💡 Action Plan (First 6 Months):"
p1.font.bold = True
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(0, 51, 102)

p2 = tf2.add_paragraph()
p2.text = "1. Global Content: Translate top global Deep Learning whitepapers to Korean."
p2.font.size = Pt(16)
p2.level = 1

p3 = tf2.add_paragraph()
p3.text = "2. Local Content: Create massive blog & news articles in Korean (Target: 100+ items)."
p3.font.size = Pt(16)
p3.level = 1

p4 = tf2.add_paragraph()
p4.text = "3. Proof: Publish 3 new local Tier-1 brand success stories."
p4.font.size = Pt(16)
p4.level = 1

p5 = tf2.add_paragraph()
p5.text = "4. Community: Launch the bi-annual 'Incrementality Summit Korea'."
p5.font.size = Pt(16)
p5.level = 1

prs.save('RTB_vs_Criteo_Awareness.pptx')
