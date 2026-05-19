#!/usr/bin/env python3
"""Generate 5-slide Leadership Plan deck — clean white, 4-color, large text."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLACK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def txt(slide, l, t, w, h, text, sz=14, clr=BLACK, bold=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = "Arial"; p.alignment = al
    return tb

def title(slide, t, sub=""):
    txt(slide, 0.7, 0.25, 11, 0.6, t, 28, BLACK, True)
    if sub: txt(slide, 0.7, 0.85, 11, 0.4, sub, 14, GRAY)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.25), Inches(1.2), Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()

def snum(slide, n):
    txt(slide, 12.4, 7.05, 0.8, 0.4, f"{n}/5", 11, GRAY, al=PP_ALIGN.RIGHT)

def table(slide, l, t, w, h, rows, cols, data, cw=None):
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = ts.table
    if cw:
        for i, v in enumerate(cw): tbl.columns[i].width = Inches(v)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12) if r > 0 else Pt(13)
                p.font.color.rgb = WHITE if r == 0 else BLACK
                p.font.bold = (r == 0); p.font.name = "Arial"
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (WHITE if r % 2 == 1 else LIGHT_BG)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
    return ts


# ============================================================
# SLIDE 1: First 6 Months & Onboarding
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s1); snum(s1, 1)
title(s1, "1. First 6 Months & Onboarding", "Challenge Identification & Risk Mitigation — 30-60-90-180 Day Plan")

txt(s1, 0.7, 1.55, 5, 0.4, "Top 3 Risks During Leadership Transition", 18, BLACK, True)
d1 = [
    ["#", "Risk", "Context", "Impact"],
    ["1", "Loss of Core Assets\n& Key Talent",
     "Failing to defend and leverage existing\ncore assets and top talent during transition",
     "Organizational momentum\nloss & pipeline collapse"],
    ["2", "Insufficient Global\nKnowledge Transfer",
     "Proven winning plays from other markets\nnot sufficiently localized to Korea",
     "Growth ceiling &\nrepeated inefficiencies"],
    ["3", "Right Person,\nRight Build Failure",
     "Recruiting misaligned talent during\naggressive team expansion",
     "Resource waste &\ndelayed market entry"],
]
table(s1, 0.7, 1.95, 11.9, 2.6, 4, 4, d1, [0.5, 2.0, 4.2, 2.8])

txt(s1, 0.7, 4.8, 5, 0.4, "Risk Mitigation — Phased Action Plan", 18, BLACK, True)
d1b = [
    ["Phase", "Period", "Key Mitigation Actions"],
    ["Phase 1: Listen", "Day 1–30",
     "1:1 Listening Tour (all 20 members) → identify flight risks and core assets early\n"
     "Shadow Week: observe Sales/AM in action without intervening\n"
     "Deliver 1 Quick Win to build credibility through action"],
    ["Phase 2: Align", "Day 31–60",
     "Build Competency Map + Flight Risk assessment → immediate retention conversations\n"
     "Design Global → Local Playbook transfer pipeline (DACH, UK, Japan best practices)\n"
     "Establish 3 weekly rituals (Mon Pipeline / Wed 1:1 / Fri Wins & Learns)"],
    ["Phase 3: Lead", "Day 61–90",
     "Declare 2026→2030 growth vision with concrete revenue targets ($42.2M → $88.7M)\n"
     "Define Right Person criteria and begin first strategic hire (Senior AE)\n"
     "Activate Delegation Framework: CM owns only Tier 1 accounts + HQ relations"],
]
table(s1, 0.7, 5.2, 11.9, 2.1, 4, 3, d1b, [1.5, 1.2, 9.2])


# ============================================================
# SLIDE 2: Integration
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2); snum(s2, 2)
title(s2, "2. Integration & Leadership Initiatives", "First 6 Months — Building Trust, Transferring Knowledge, Enabling Growth")

sections = [
    ("Talent Retention & Leverage (Risk #1)", [
        "Map every team member's core competencies and flight risk level by Day 60",
        "Immediate 1:1 career conversations for high-risk seniors — present clear growth paths",
        "Convert tacit knowledge (client histories, deal playbooks) into shared team assets",
        "Grant clear ownership and decision-making authority → \"my role is growing here\"",
    ]),
    ("Global Knowledge Transfer (Risk #2)", [
        "Audit Top 3 winning plays from DACH, UK, France, Japan markets",
        "Apply Korea Localization Filter — what works, what doesn't in agency-driven market",
        "Launch monthly \"Global → Local\" sessions (bidirectional: Korea exports wins too)",
        "Immediately apply Criteo win-back playbooks from Europe to Commerce Squad",
    ]),
    ("Right Person, Right Build (Risk #3)", [
        "Use 60-day weekly cadence to observe org strengths and gaps before any hiring",
        "Define Right Person criteria: ① DL/programmatic fluency ② A/B test sales ③ culture fit",
        "First hire: 1 Senior AE with mega-account closing experience ($70K+/mo deals)",
        "Principle: \"Fill the seat correctly, not quickly\" — 1 wrong hire delays growth 6 months",
    ]),
]
for i, (heading, bullets) in enumerate(sections):
    x = 0.7 + i * 4.2
    txt(s2, x, 1.55, 3.8, 0.4, heading, 16, BLUE, True)
    for j, b in enumerate(bullets):
        txt(s2, x, 2.1 + j * 0.55, 3.9, 0.5, f"•  {b}", 12, BLACK)

# Weekly cadence table
txt(s2, 0.7, 4.5, 5, 0.4, "Weekly Operating Cadence (3 Rituals Only)", 16, BLACK, True)
dc = [
    ["Day", "Meeting", "Participants", "Duration", "Purpose"],
    ["Mon", "Pipeline Review", "All Sales + AM", "30 min", "Pipeline status + competency gap observation"],
    ["Wed", "1:1 Coaching", "Individual (rotating)", "30 min", "Deal coaching + growth direction (behavior-based)"],
    ["Fri", "Wins & Learns", "Entire team (voluntary)", "20 min", "1 success + 1 failure lesson → learning culture"],
]
table(s2, 0.7, 4.9, 11.9, 1.4, 4, 5, dc, [0.7, 1.8, 2.0, 1.2, 6.2])

txt(s2, 0.7, 6.6, 11.9, 0.5,
    "Core Principle:  \"Don't manage — coach. Don't dictate — bridge. Don't assume — listen.\"",
    16, BLUE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Structural Evolution
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s3); snum(s3, 3)
title(s3, "3. Structural Evolution", "Short-Term Team Structure (6–12 Months) — Vertical Squads by Industry")

txt(s3, 0.7, 1.55, 10, 0.4, "Vertical Squad Organization (Day 91–180)", 18, BLACK, True)
txt(s3, 0.7, 1.95, 10, 0.35,
    "Reorganize existing 20-person team by redefining roles to maximize core talent strengths — not a large-scale overhaul.",
    13, GRAY)

ds = [
    ["Squad", "Composition", "Target Accounts", "2026 Mission", "Market Rationale"],
    ["🛒 Commerce\nSquad",
     "Sr. Sales 1\nAM 2, TAM 1",
     "Coupang, Musinsa,\nOlive Young, SSG",
     "Onboard 2 new\nTier 1 clients",
     "E-Commerce = $3.2B–$3.5B (45–50%\nof digital ad spend). #1 Criteo win-back path."],
    ["🎮 Gaming &\nApp Squad",
     "Sales 1\nAM 1, TAM 1",
     "Nexon, Krafton,\nNetmarble",
     "1 Tier 1 A/B test\nsuccess → conversion",
     "Gaming = $1.1B–$1.4B budget.\nPredict LTV to beat Moloco/AppLovin."],
    ["🌏 K-Brand &\nGrowth Squad",
     "Sales 1\nAM 1",
     "APR, Samyang,\nYanolja, Toss",
     "Build 2 global D2C\npipelines",
     "Hybrid K-Brand = $1.4B–$1.8B.\nGlobal CBEC = explosive growth driver."],
]
table(s3, 0.7, 2.45, 11.9, 3.0, 4, 5, ds, [1.3, 1.3, 1.8, 1.8, 5.7])

txt(s3, 0.7, 5.7, 5, 0.4, "Lean Hiring Roadmap", 18, BLACK, True)
dh = [
    ["Period", "Headcount", "Key Hire", "Purpose"],
    ["2026 (Now → 6)", "5 → 6", "1 Senior AE (mega-account closer)", "Breakthrough for H2 44 new launches, Tier 1 deals"],
    ["2027–28 (8 → 10)", "+4", "2 AEs + 2 Campaign Managers", "Match 25+ launches/quarter pace, K-Brand D2C dedicated"],
    ["2029–30 (11 → 12)", "+2", "1 Data Scientist + 1 Sr. AE", "Complete 495 active campaign operating system"],
]
table(s3, 0.7, 6.1, 11.9, 1.2, 4, 4, dh, [1.8, 1.2, 3.5, 5.4])


# ============================================================
# SLIDE 4: Team Collaboration
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s4); snum(s4, 4)
title(s4, "4. Team Collaboration", "Sales & AM Alignment, Motivation, and Driving Results")

# Two rituals
txt(s4, 0.7, 1.55, 6, 0.4, "2 Collaboration Rituals (Simple = Executable)", 18, BLACK, True)

cols_data = [
    ("① Pre-flight Handoff", [
        "Sales prepares a 1-page Campaign Brief before every AM handoff",
        "Includes: client KPIs, budget/terms, competitive landscape, decision-maker map",
        "Korea-specific: must specify whether final decision-maker is advertiser or agency",
        "This single line dramatically improves AM communication efficiency",
    ]),
    ("② Monthly Win/Loss Review", [
        "30 min/month. Sales + AM jointly analyze 1 won deal and 1 lost deal",
        "Rule: No blame. Only discuss \"how can we win next time?\"",
        "CM shares own failure first → builds psychological safety",
        "Korean culture: 'sharing failures' conflicts with face-saving → CM leads by example",
    ]),
]
for i, (heading, bullets) in enumerate(cols_data):
    x = 0.7 + i * 6.2
    txt(s4, x, 2.0, 5.8, 0.4, heading, 16, BLUE, True)
    for j, b in enumerate(bullets):
        txt(s4, x, 2.5 + j * 0.45, 5.8, 0.4, f"•  {b}", 12, BLACK)

# Motivation
txt(s4, 0.7, 4.7, 6, 0.4, "Motivation System", 18, BLACK, True)

mot = [
    ("💰  Monetary", [
        "Team-wide bonus on quarterly target achievement (separate from individual commission)",
        "Team-level rewards prevent silos and reinforce collaboration culture",
    ]),
    ("🌟  Non-Monetary", [
        "Quarterly Top Performer: Warsaw HQ R&D Center visit + global workshop participation",
        "Aligned with RTB House values: Kindness, Curiosity, Excellence, Ownership",
        "Non-monetary growth opportunities are the strongest retention lever for Korean talent",
    ]),
]
for i, (heading, bullets) in enumerate(mot):
    x = 0.7 + i * 6.2
    txt(s4, x, 5.15, 5.8, 0.35, heading, 16, BLUE, True)
    for j, b in enumerate(bullets):
        txt(s4, x, 5.55 + j * 0.42, 5.8, 0.4, f"•  {b}", 12, BLACK)


# ============================================================
# SLIDE 5: Performance Measurement
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s5); snum(s5, 5)
title(s5, "5. Performance Measurement", "KPIs, OKRs & Revenue Targets — 70% Quantitative + 30% Qualitative")

txt(s5, 0.7, 1.55, 5, 0.35, "A. Sales KPIs & OKRs", 17, BLACK, True)
dsa = [
    ["Type", "Metric", "Frequency", "Year 1 Target"],
    ["Quantitative", "New Bookings (MRR)", "Monthly", "$150K/mo → $800K by Q4"],
    ["Quantitative", "New Campaign Launches", "Quarterly", "10+ per quarter (min $7K/mo test)"],
    ["Quantitative", "Pipeline Coverage", "Weekly", "3x quarterly target maintained"],
    ["Qualitative OKR", "O: RTB House = \"#1 DL Performance\" in Korea e-commerce", "", ""],
    ["", "KR1: C-level DL seminar 1x/quarter", "Quarterly", ""],
    ["", "KR2: 2+ Korean case studies published", "Semi-annual", ""],
]
table(s5, 0.7, 1.9, 11.9, 2.0, 7, 4, dsa, [1.3, 4.5, 1.3, 4.8])

txt(s5, 0.7, 4.1, 5, 0.35, "B. Account Management KPIs & OKRs", 17, BLACK, True)
dam = [
    ["Type", "Metric", "Frequency", "Year 1 Target"],
    ["Quantitative", "Scale-up Rate (test → regular)", "Quarterly", "60%+ conversion rate"],
    ["Quantitative", "Campaign Churn Rate", "Monthly", "Below 5% monthly"],
    ["Quantitative", "Time-to-Value", "Per campaign", "Target ROAS within 21 days"],
    ["Qualitative OKR", "O: Become client's strategic business partner", "", ""],
    ["", "KR1: 100% QBR completion for major clients", "Quarterly", ""],
    ["", "KR2: Programmatic & cookieless training 2x/yr", "Semi-annual", ""],
]
table(s5, 0.7, 4.45, 11.9, 2.0, 7, 4, dam, [1.3, 4.5, 1.3, 4.8])

txt(s5, 0.7, 6.65, 5, 0.35, "C. Country Manager Metrics (HQ Reporting)", 17, BLACK, True)
dcm = [
    ["Metric", "Year 1 Target"],
    ["Revenue Growth (YoY)", "+30% (foundation year)"],
    ["Team Attrition Rate", "Below 10% (vs industry avg 20%+)"],
    ["New Vertical Expansion", "First revenue in 2+ new verticals"],
    ["Internal NPS", "8.0/10+ (semi-annual anonymous survey)"],
]
table(s5, 6.5, 6.65, 6.1, 0.75, 5, 2, dcm, [2.8, 3.3])

# Save
out = '/Users/wook/WookAi/Booklog/RTB_House_Leadership_Plan_Deck.pptx'
prs.save(out)
print(f"✅ Saved: {out}")
