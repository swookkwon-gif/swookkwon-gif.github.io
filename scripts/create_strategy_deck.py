#!/usr/bin/env python3
"""RTB House Korea Strategy Deck — Clean, readable, 4-color, USD-only."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# === 4 Colors Only ===
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)  # table alt row (near-white, not a "new" color)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE


def add_text(slide, left, top, width, height, text, size=14,
             color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tb


def add_title(slide, title, subtitle=""):
    add_text(slide, 0.7, 0.3, 11, 0.7, title, 32, BLACK, bold=True)
    if subtitle:
        add_text(slide, 0.7, 0.95, 11, 0.5, subtitle, 16, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.7), Inches(1.4), Inches(1.2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13) if r > 0 else Pt(14)
                p.font.color.rgb = WHITE if r == 0 else BLACK
                p.font.bold = (r == 0)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = BLUE
            elif r % 2 == 1:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LIGHT_BG
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
    return ts


def slide_num(slide, n):
    add_text(slide, 12.4, 7.05, 0.8, 0.4, f"{n} / 5", 11, GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: High-Level Strategy
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
add_title(s1, "1. High-Level Strategy", "Focus Industries & Rationale")
slide_num(s1, 1)

d1 = [
    ["#", "Sector", "Budget", "Why Prioritized"],
    ["1", "E-Commerce & Retail",
     "$2.1B – $2.5B",
     "45–50% of Korea's digital ad spend. Deep learning dynamic banner\noptimization in mega-catalog (100K+ SKU) environments. #1 Criteo win-back target."],
    ["2", "K-Brands Global D2C",
     "$1.4B – $1.8B",
     "Samsung, APR/Medicube, Samyang. Massive budgets shifting to\nglobal cross-border (CBEC) purchase conversions. Fastest revenue growth driver."],
    ["3", "Gaming & App",
     "$1.1B – $1.4B",
     "CPI/LTV optimization = survival for publishers. Predictive LTV modeling\nproves efficiency vs Moloco & AppLovin in UA campaigns."],
    ["4", "Legacy Giants",
     "$1.8B – $2.1B",
     "Hyundai, Telcos, Finance. Branding-focused KPIs. Mid-long term pipeline\nvia ContextAI (cookieless), Brand Safety, and in-house agency partnerships."],
]
add_table(s1, 0.7, 1.7, 11.9, 4.5, 5, 4, d1, col_widths=[0.5, 2.5, 1.8, 7.1])

add_text(s1, 0.7, 6.5, 11.9, 0.5,
         "Core Message: \"Criteo competes with massive data. RTB House competes with smarter algorithms.\"",
         15, BLUE, bold=True)


# ============================================================
# SLIDE 2: Top Prospective Clients
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
add_title(s2, "2. Top Prospective Clients", "Tier 1 Targets (Annual Budget $3.5M+) & Deal Strategy")
slide_num(s2, 2)

d2 = [
    ["Target", "Sector", "Ad Spend", "Current", "Attack Strategy"],
    ["Coupang / SSG /\nGmarket / 11st",
     "E-Commerce", "$947M+\n(10-K)", "Criteo",
     "30-day free A/B test. Measure iROAS vs Criteo\nat 50:50 traffic split in mega-catalog."],
    ["Samsung / LG\nElectronics",
     "Global Tech", "$1B+\n(Global)", "Multi-DSP",
     "Global D2C retargeting for electronics malls.\nTarget direct traffic → purchase ROAS."],
    ["APR / Olive Young /\nAmorepacific",
     "K-Beauty", "$100M+", "Meta,\nCriteo",
     "Global targeting for beauty device NA expansion\nand cross-channel purchase induction."],
    ["Nexon / Krafton /\nNetmarble",
     "Gaming", "$100M+\n/company", "Moloco,\nAppLovin",
     "UA for global titles. Prove LTV prediction\nsuperiority via predictive deep learning."],
]
add_table(s2, 0.7, 1.7, 11.9, 3.3, 5, 5, d2, col_widths=[2.2, 1.3, 1.2, 1.0, 6.2])

# 3-Step Framework
add_text(s2, 0.7, 5.3, 8, 0.4, "3-Step Deal Framework", 20, BLACK, bold=True)

steps_text = [
    ("① PROVE IT", "30-day free A/B test\n50:50 split vs incumbent DSP\nMeasure Incremental Conversions"),
    ("② SCALE IT", "On success, migrate budget\n30% → 50% → 70% to RTB House"),
    ("③ LOCK IT", "Quarterly iROAS reporting\nvia Ghost Ads methodology\nRaise Switching Costs"),
]
for i, (t, d) in enumerate(steps_text):
    x = 0.7 + i * 4.1
    add_text(s2, x, 5.8, 3.5, 0.4, t, 18, BLUE, bold=True)
    add_text(s2, x, 6.25, 3.5, 1.0, d, 14, GRAY)


# ============================================================
# SLIDE 3: Campaign Projections
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
add_title(s3, "3. Campaign Projections", "Quarterly Launch Targets (2026–2030)")
slide_num(s3, 3)

d3 = [
    ["Year", "Q1", "Q2", "Q3", "Q4", "Annual\nLaunches", "Active\nCampaigns", "Revenue"],
    ["2026 (Foundation)", "—", "—", "20", "24", "44", "235", "$42.2M"],
    ["2027 (Scale-up)", "24", "25", "26", "26", "101", "308", "$55.4M"],
    ["2028 (Expansion)", "24", "25", "26", "26", "101", "372", "$66.8M"],
    ["2029 (Acceleration)", "25", "26", "26", "27", "104", "431", "$78.4M"],
    ["2030 (Reclaim #1)", "28", "29", "29", "30", "116", "495", "$94.3M"],
]
add_table(s3, 0.7, 1.7, 11.9, 3.6, 6, 8, d3,
          col_widths=[2.2, 0.9, 0.9, 0.9, 0.9, 1.2, 1.4, 1.5])

# Key assumptions
add_text(s3, 0.7, 5.6, 8, 0.4, "Key Assumptions", 20, BLACK, bold=True)

assumptions = [
    "Blended Revenue per Campaign:  $179K/yr  (90% standard $105K + 10% mega $844K)",
    "Annual Churn Rate:  12% fixed  (B2B SaaS benchmark, conservative)",
    "2030 Target:  495 active campaigns  — only 60% of Criteo's volume (~840) needed to overtake",
]
for i, text in enumerate(assumptions):
    y = 6.1 + i * 0.38
    add_text(s3, 0.9, y, 11.5, 0.35, f"•  {text}", 14, GRAY)


# ============================================================
# SLIDE 4: Revenue Metrics
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
add_title(s4, "4. Revenue Metrics", "Per-Salesperson Projections — Lean Organization Model")
slide_num(s4, 4)

d4 = [
    ["Year", "Headcount", "New Booking\n/ Person", "Gross\nNew", "Churn\n(12%)", "Net\nGrowth", "Revenue"],
    ["2026", "6", "$1.76M", "$10.6M", "–$4.4M", "+$6.2M", "$42.6M"],
    ["2027", "8", "$1.79M", "$14.3M", "–$5.1M", "+$9.3M", "$51.9M"],
    ["2028", "10", "$1.86M", "$18.6M", "–$6.2M", "+$12.5M", "$64.4M"],
    ["2029", "11", "$1.97M", "$21.7M", "–$7.7M", "+$14.0M", "$78.4M"],
    ["2030", "12", "$2.11M", "$25.3M", "–$9.4M", "+$16.0M", "$94.3M"],
]
add_table(s4, 0.7, 1.7, 11.9, 3.5, 6, 7, d4,
          col_widths=[1.2, 1.2, 1.5, 1.5, 1.5, 1.5, 1.5])

# 2026 quarterly breakdown
add_text(s4, 0.7, 5.5, 10, 0.4, "2026 Per-Salesperson Quarterly Detail (6 AEs)", 20, BLACK, bold=True)

dq = [
    ["Metric", "Q2 (Ramp-up)", "Q3", "Q4", "H2 Total"],
    ["New Campaigns / Person", "—", "3–4", "4", "~7"],
    ["New Booking / Person", "—", "$774K", "$985K", "$1.76M"],
    ["Team New Launches", "—", "20", "24", "44"],
    ["Team New Revenue", "—", "$4.6M", "$6.0M", "$10.6M"],
]
add_table(s4, 0.7, 6.0, 11.9, 1.4, 5, 5, dq, col_widths=[2.8, 2.1, 2.1, 2.1, 2.8])


# ============================================================
# SLIDE 5: Operational Activities
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
add_title(s5, "5. Operational Activities", "Building a Successful Foundation Beyond Sales")
slide_num(s5, 5)

sections = [
    ("Lean Team & Employer Branding", [
        "Position as 'Europe's #1 Deep Learning AI Tech Company'",
        "2026: +1 Senior AE (mega-account closer, $70K+/mo deals)",
        "2027–28: +2 AEs, +2 CMs for 25+ launches/quarter",
        "Quarterly 'Deep Learning in AdTech' Tech Talks",
        "Seoul National Univ. / KAIST AI lab internship pipeline",
    ]),
    ("Innovative Solution Selling", [
        "ContextAI: Cookieless targeting via DL+NLP context analysis",
        "  — RTB House OBTD framework adopted by Google Privacy Sandbox",
        "Video Retargeting: Auto 15s dynamic video ads (+30–50% CTR)",
        "Predictive LTV Modeling: 14-day LTV for gaming UA campaigns",
        "  — Prove superiority vs Moloco & AppLovin",
    ]),
    ("Market Education & Client Success", [
        "Incrementality Whitepapers: Shift market from last-click ROAS",
        "  to Ghost Ads-based iROAS measurement",
        "Quarterly QBRs with iROAS reporting for all major clients",
        "Goal: Lower annual churn from 12% → below 10%",
        "Bi-annual 'Incrementality Summit Korea' for pipeline building",
    ]),
]

for i, (title, bullets) in enumerate(sections):
    x = 0.7 + i * 4.2
    add_text(s5, x, 1.7, 3.8, 0.4, title, 18, BLUE, bold=True)
    for j, b in enumerate(bullets):
        y = 2.25 + j * 0.42
        indent = b.startswith("  ")
        add_text(s5, x + (0.3 if indent else 0), y, 3.8, 0.4,
                 ("→ " + b.strip()) if indent else ("•  " + b), 13,
                 GRAY if indent else BLACK)

# Bottom goal
add_text(s5, 0.7, 6.8, 11.9, 0.5,
         "2030 Goal:  Overtake Criteo Korea ($88.7M) → Reach $94.3M with 12 elite salespeople",
         18, BLUE, bold=True, align=PP_ALIGN.CENTER)


# === Save ===
out = "/Users/wook/WookAi/Booklog/RTB_House_Korea_Strategy_Deck.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
