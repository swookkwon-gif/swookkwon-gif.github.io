#!/usr/bin/env python3
"""Generate 1-slide PPTX: RTB House 2030 KR Dominance Plan with bar chart."""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# === Data ===
years =    [2019, 2020, 2021, 2022, 2023, 2024, 2025, 2026, 2027, 2028, 2029, 2030]
revenue =  [0.8,  2.9,  7.2,  28.0, 35.8, 47.3, 35.5, 39.7, 48.6, 59.4, 72.6, 88.7]
labels =   ['2019','2020','2021','2022','2023','2024','2025','2026(e)','2027(e)','2028(e)','2029(e)','2030(e)']
actual =   [True]*7 + [False]*5

# === Colors (4-color palette) ===
BLUE = '#2563EB'
BLUE_LIGHT = '#93C5FD'
BLACK = '#1A1A2E'
GRAY = '#6B7280'

# === Generate Chart ===
fig, ax = plt.subplots(figsize=(11, 4.8))
fig.patch.set_facecolor('white')
ax.set_facecolor('white')

colors = [BLUE if a else BLUE_LIGHT for a in actual]
bars = ax.bar(range(len(years)), revenue, color=colors, width=0.65, edgecolor='white', linewidth=0.5)

# Value labels on bars
for i, (bar, val) in enumerate(zip(bars, revenue)):
    y_pos = bar.get_height() + 1.2
    fontsize = 11 if val >= 10 else 10
    ax.text(bar.get_x() + bar.get_width()/2, y_pos, f'${val:.1f}M',
            ha='center', va='bottom', fontsize=fontsize, fontweight='bold',
            color=BLACK)



# Styling
ax.set_xticks(range(len(labels)))
ax.set_xticklabels(labels, fontsize=11, color=BLACK)
ax.set_ylabel('Revenue ($M)', fontsize=12, color=BLACK, labelpad=10)
ax.set_ylim(0, 105)
ax.yaxis.set_major_formatter(mticker.FormatStrFormatter('$%.0fM'))
ax.tick_params(axis='y', labelsize=10, colors=GRAY)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.spines['left'].set_color('#E5E7EB')
ax.spines['bottom'].set_color('#E5E7EB')
ax.yaxis.grid(True, color='#F3F4F6', linewidth=0.8)
ax.set_axisbelow(True)

# Legend
from matplotlib.patches import Patch
legend_elements = [Patch(facecolor=BLUE, label='Actual'),
                   Patch(facecolor=BLUE_LIGHT, label='Estimated')]
ax.legend(handles=legend_elements, loc='upper left', fontsize=10,
          frameon=False, labelcolor=BLACK)

plt.tight_layout(pad=1.0)
chart_path = '/Users/wook/WookAi/Booklog/scripts/rtb_chart.png'
plt.savefig(chart_path, dpi=200, bbox_inches='tight', facecolor='white')
plt.close()
print(f"✅ Chart saved: {chart_path}")

# === Build PPTX ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Title
tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(11), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Summary: RTB House 2030 KR Dominance Plan"
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.font.name = "Arial"

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(0.9), Inches(1.2), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
line.line.fill.background()

# Subtitle
tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(11), Inches(0.4))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "RTB House Korea Revenue Trajectory & 2030 Growth Forecast — Overtaking Criteo Korea"
p2.font.size = Pt(15)
p2.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
p2.font.name = "Arial"

# Chart image
slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.55), Inches(12.3), Inches(4.7))

# Key Assumptions box
tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(12), Inches(0.4))
tf3 = tb3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Key Assumptions"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p3.font.name = "Arial"

assumptions = [
    "•  2025 baseline: Estimated at 40% of Criteo Korea's DART-disclosed revenue ($88.7M × 40% = $35.5M)",
    "•  2030 target: Grow to match Criteo Korea's 2025 revenue level ($88.7M) — reclaiming market #1",
    "•  CAGR 2025→2030: ~20%  |  Churn rate: 12% annual (B2B SaaS benchmark)  |  Team: 6 → 12 elite salespeople",
]
for i, text in enumerate(assumptions):
    tb4 = slide.shapes.add_textbox(Inches(0.9), Inches(6.8 + i * 0.28), Inches(11.5), Inches(0.3))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = text
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p4.font.name = "Arial"

# Save
out = '/Users/wook/WookAi/Booklog/RTB_House_2030_Dominance_Plan.pptx'
prs.save(out)
print(f"✅ PPTX saved: {out}")
